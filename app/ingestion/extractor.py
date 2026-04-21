"""
Ingestion Pipeline – Dokumenten-Extraktion.
Unterstützt: PDF, DOCX, PPTX, XLSX, MD, TXT, XML, RTF.
Extrahiert Text, erkennt Bilder & Tabellen.
"""
import io
import logging
import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ExtractionResult:
    """Ergebnis der Dokumenten-Extraktion."""
    text: str = ""
    page_count: int = 0
    has_images: bool = False
    has_tables: bool = False
    metadata: dict = field(default_factory=dict)
    errors: list[str] = field(default_factory=list)


class DocumentExtractor:
    """
    Multi-Format-Extraktor.
    Lokal, keine Cloud-APIs – alle Bibliotheken laufen on-premise.
    """

    def extract(self, file_path: Path) -> ExtractionResult:
        """Hauptmethode: Format erkennen und extrahieren."""
        suffix = file_path.suffix.lower()
        extractors = {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".doc": self._extract_docx,  # Fallback
            ".pptx": self._extract_pptx,
            ".ppt": self._extract_pptx,
            ".xlsx": self._extract_xlsx,
            ".xls": self._extract_xlsx,
            ".md": self._extract_text,
            ".txt": self._extract_text,
            ".xml": self._extract_xml,
            ".rtf": self._extract_rtf,
        }

        extractor = extractors.get(suffix)
        if not extractor:
            return ExtractionResult(errors=[f"Format nicht unterstützt: {suffix}"])

        try:
            result = extractor(file_path)
            # Text bereinigen
            result.text = self._clean_text(result.text)
            return result
        except Exception as e:
            logger.error(f"Extraktion fehlgeschlagen für {file_path}: {e}")
            return ExtractionResult(errors=[str(e)])

    # ─── PDF ───

    def _extract_pdf(self, path: Path) -> ExtractionResult:
        import pdfplumber

        result = ExtractionResult()
        text_parts = []

        with pdfplumber.open(str(path)) as pdf:
            result.page_count = len(pdf.pages)
            result.metadata = {
                k: v for k, v in (pdf.metadata or {}).items()
                if isinstance(v, (str, int, float))
            }

            for page in pdf.pages:
                # Text
                page_text = page.extract_text() or ""
                text_parts.append(page_text)

                # Tabellen erkennen
                tables = page.extract_tables()
                if tables:
                    result.has_tables = True
                    for table in tables:
                        if table:
                            for row in table:
                                cells = [str(c) if c else "" for c in row]
                                text_parts.append(" | ".join(cells))

                # Bilder erkennen
                if page.images:
                    result.has_images = True

        result.text = "\n\n".join(text_parts)
        return result

    # ─── DOCX ───

    def _extract_docx(self, path: Path) -> ExtractionResult:
        from docx import Document as DocxDocument

        result = ExtractionResult()
        doc = DocxDocument(str(path))
        text_parts = []

        # Kern-Eigenschaften
        if doc.core_properties:
            props = doc.core_properties
            result.metadata = {
                "author": props.author or "",
                "title": props.title or "",
                "subject": props.subject or "",
                "created": str(props.created) if props.created else "",
            }

        # Absätze
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Tabellen
        for table in doc.tables:
            result.has_tables = True
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                text_parts.append(" | ".join(cells))

        # Bilder erkennen (via Inline-Shapes)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                result.has_images = True
                break

        result.text = "\n\n".join(text_parts)
        result.page_count = max(1, len(text_parts) // 40)  # Schätzung
        return result

    # ─── PPTX ───

    def _extract_pptx(self, path: Path) -> ExtractionResult:
        from pptx import Presentation

        result = ExtractionResult()
        prs = Presentation(str(path))
        text_parts = []
        result.page_count = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- Folie {slide_idx} ---"]
            for shape in slide.shapes:
                # Text
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

                # Tabellen
                if shape.has_table:
                    result.has_tables = True
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(cells))

                # Bilder
                if hasattr(shape, "image"):
                    result.has_images = True

            text_parts.append("\n".join(slide_texts))

        result.text = "\n\n".join(text_parts)
        return result

    # ─── XLSX ───

    def _extract_xlsx(self, path: Path) -> ExtractionResult:
        from openpyxl import load_workbook

        result = ExtractionResult()
        result.has_tables = True
        text_parts = []

        wb = load_workbook(str(path), read_only=True, data_only=True)
        result.page_count = len(wb.sheetnames)

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"--- Blatt: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    text_parts.append(" | ".join(cells))

        wb.close()
        result.text = "\n\n".join(text_parts)
        return result

    # ─── Markdown / Text ───

    def _extract_text(self, path: Path) -> ExtractionResult:
        result = ExtractionResult()
        text = path.read_text(encoding="utf-8", errors="replace")
        result.text = text
        result.page_count = max(1, text.count("\n") // 50)

        # Markdown-Tabellen erkennen
        if re.search(r"\|.*\|.*\|", text):
            result.has_tables = True

        # Markdown-Bilder erkennen
        if re.search(r"!\[.*?\]\(.*?\)", text):
            result.has_images = True

        return result

    # ─── XML ───

    def _extract_xml(self, path: Path) -> ExtractionResult:
        from bs4 import BeautifulSoup

        result = ExtractionResult()
        raw = path.read_text(encoding="utf-8", errors="replace")
        soup = BeautifulSoup(raw, "lxml-xml")
        result.text = soup.get_text(separator="\n", strip=True)
        result.page_count = 1
        result.metadata = {"root_tag": soup.find().name if soup.find() else "unknown"}
        return result

    # ─── RTF ───

    def _extract_rtf(self, path: Path) -> ExtractionResult:
        from striprtf.striprtf import rtf_to_text

        result = ExtractionResult()
        raw = path.read_text(encoding="utf-8", errors="replace")
        result.text = rtf_to_text(raw)
        result.page_count = max(1, result.text.count("\n") // 50)
        return result

    # ─── Hilfsmethoden ───

    @staticmethod
    def _clean_text(text: str) -> str:
        """Text bereinigen: überflüssige Leerzeichen und Zeilenumbrüche entfernen."""
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r" +\n", "\n", text)
        return text.strip()


# Singleton
extractor = DocumentExtractor()
