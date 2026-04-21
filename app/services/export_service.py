"""
Export-Service – Echte Datei-Generierung (PPTX, DOCX, PDF).

Erzeugt reale Office-Dateien aus Store-Inhalten:
  - PPTX: python-pptx mit Komm.ONE CI
  - DOCX: python-docx mit Formatierung
  - PDF:  reportlab (Zusammenfassungen, Live-View)

Alle Exporte basieren ausschliesslich auf dem jeweiligen Store.
"""
import logging
import io
from pathlib import Path
from typing import Optional

from app.core.config import settings
from app.services.intelligence import (
    generate_summary, extract_key_takeaways, fuse_knowledge, distill_facts,
)

logger = logging.getLogger(__name__)

# Komm.ONE CI Farben (RGB-Tupel)
CI_MIDNIGHT = (0x00, 0x3A, 0x40)
CI_LAGOON = (0x00, 0xB2, 0xA9)
CI_AMARILLO = (0xF1, 0xC4, 0x00)
CI_WHITE = (0xFF, 0xFF, 0xFF)
CI_GRAY = (0x6C, 0x75, 0x7D)
CI_LIGHT = (0xF2, 0xF5, 0xF5)


def export_pptx(
    store_name: str,
    store_type: str,
    documents: list,
    params: dict = None,
) -> bytes:
    """
    PowerPoint-Praesentation aus Store-Inhalten generieren.
    Gibt die PPTX-Datei als Bytes zurueck.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    params = params or {}
    title = params.get("title", f"Bericht: {store_name}")
    focus = params.get("focus", "")

    texts = [d.get("content", "") or (d.content_text if hasattr(d, "content_text") else "") for d in documents]
    texts = [t for t in texts if t]
    summary = generate_summary(texts, max_sentences=6)
    takeaways = extract_key_takeaways(texts, max_items=8)
    facts = distill_facts(texts)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(layout_idx=6):
        return prs.slides.add_slide(prs.slide_layouts[layout_idx])

    def add_text(slide, left, top, width, height, text, size=18, bold=False, color=CI_MIDNIGHT, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.alignment = align
        return tf

    def add_bg(slide, color=CI_WHITE):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    # ── Folie 1: Titelfolie ──
    s = add_slide()
    add_bg(s, CI_MIDNIGHT)
    add_text(s, 1, 1.5, 11, 1.5, title, size=36, bold=True, color=CI_WHITE, align=PP_ALIGN.CENTER)
    type_label = "Akte" if store_type == "akte" else "WissensDB"
    add_text(s, 1, 3.2, 11, 0.8, f"{type_label}: {store_name}", size=20, color=CI_LAGOON, align=PP_ALIGN.CENTER)
    add_text(s, 1, 4.5, 11, 0.6, f"Basierend auf {len(documents)} Dokumenten | Agentischer Document Store",
             size=14, color=CI_GRAY, align=PP_ALIGN.CENTER)

    # ── Folie 2: Zusammenfassung ──
    s = add_slide()
    add_bg(s)
    add_text(s, 0.8, 0.5, 11, 0.7, "Zusammenfassung", size=28, bold=True, color=CI_LAGOON)
    add_text(s, 0.8, 1.5, 11, 5.0, summary[:800], size=16, color=CI_MIDNIGHT)

    # ── Folie 3: Kernfakten ──
    s = add_slide()
    add_bg(s)
    add_text(s, 0.8, 0.5, 11, 0.7, "Kernfakten & Takeaways", size=28, bold=True, color=CI_AMARILLO)
    y = 1.5
    for tk in takeaways[:8]:
        val = tk.get("takeaway", tk.get("entity", ""))
        cnt = tk.get("count", 0)
        add_text(s, 1.0, y, 10, 0.4, f"  {val}  ({cnt}x)", size=16, color=CI_MIDNIGHT)
        y += 0.45

    # ── Folie 4: Destillierte Fakten ──
    if facts:
        s = add_slide()
        add_bg(s)
        add_text(s, 0.8, 0.5, 11, 0.7, "Extrahierte Fakten", size=28, bold=True, color=CI_LAGOON)
        y = 1.5
        for f in facts[:8]:
            text = f.get("text", "")[:120]
            ftype = f.get("type", "fakt")
            prefix = {"zahl": "📊", "datum": "📅", "fakt": "📌"}.get(ftype, "•")
            add_text(s, 1.0, y, 10, 0.5, f"{prefix} {text}", size=14, color=CI_MIDNIGHT)
            y += 0.5

    # ── Folie 5+: Dokument-Uebersicht ──
    s = add_slide()
    add_bg(s)
    add_text(s, 0.8, 0.5, 11, 0.7, "Dokumenten-Uebersicht", size=28, bold=True, color=CI_LAGOON)
    y = 1.5
    for doc in documents[:10]:
        doc_title = doc.get("title", "") or (doc.title if hasattr(doc, "title") else "")
        doc_type = doc.get("file_type", "") or (doc.file_type if hasattr(doc, "file_type") else "")
        pages = doc.get("page_count", 0) or (doc.page_count if hasattr(doc, "page_count") else 0)
        add_text(s, 1.0, y, 10, 0.4, f"  {doc_title}  ({doc_type.upper()}, {pages} S.)", size=14, color=CI_MIDNIGHT)
        y += 0.4

    # ── Letzte Folie: Quellen ──
    s = add_slide()
    add_bg(s, CI_MIDNIGHT)
    add_text(s, 1, 2.5, 11, 1.0, "Quellennachweis", size=28, bold=True, color=CI_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, 1, 3.8, 11, 0.8,
             f"Alle Inhalte stammen ausschliesslich aus der {type_label} '{store_name}'.",
             size=16, color=CI_LAGOON, align=PP_ALIGN.CENTER)
    add_text(s, 1, 4.8, 11, 0.6, "Agentischer Document Store | Komm.ONE | On-Premise | DSGVO-konform",
             size=12, color=CI_GRAY, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info(f"PPTX Export: {len(prs.slides)} Folien fuer '{store_name}'")
    return buf.getvalue()


def export_docx(
    store_name: str,
    store_type: str,
    documents: list,
    entities_data: dict = None,
    params: dict = None,
) -> bytes:
    """
    Word-Dokument aus Store-Inhalten generieren.
    """
    from docx import Document as DocxDocument
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    params = params or {}
    doc_title = params.get("title", f"Bericht: {store_name}")
    sections_param = params.get("sections", "Zusammenfassung,Kernfakten,Entitaeten,Dokumentenuebersicht,Quellen")
    requested_sections = [s.strip() for s in sections_param.split(",")]

    texts = [d.get("content", "") or (d.content_text if hasattr(d, "content_text") else "") for d in documents]
    texts = [t for t in texts if t]
    summary = generate_summary(texts, max_sentences=10)
    takeaways = extract_key_takeaways(texts, max_items=12)
    facts = distill_facts(texts)

    doc = DocxDocument()

    # Styles
    style = doc.styles["Normal"]
    style.font.name = "Arial"
    style.font.size = Pt(11)

    # Titel
    h = doc.add_heading(doc_title, level=0)
    for run in h.runs:
        run.font.color.rgb = RGBColor(*CI_MIDNIGHT)

    type_label = "Akte" if store_type == "akte" else "WissensDB"
    p = doc.add_paragraph(f"{type_label}: {store_name} | {len(documents)} Dokumente")
    p.style.font.color.rgb = RGBColor(*CI_GRAY)
    doc.add_paragraph("")

    for section in requested_sections:
        sl = section.lower()

        if "zusammenfassung" in sl or "summary" in sl or "einleitung" in sl:
            doc.add_heading("Zusammenfassung", level=1)
            doc.add_paragraph(summary)

        elif "kernfakt" in sl or "takeaway" in sl:
            doc.add_heading("Kernfakten", level=1)
            for tk in takeaways:
                val = tk.get("takeaway", tk.get("entity", ""))
                cnt = tk.get("count", 0)
                doc.add_paragraph(f"{val} ({cnt}x)", style="List Bullet")

        elif "fakt" in sl or "destill" in sl:
            doc.add_heading("Extrahierte Fakten", level=1)
            for f in facts[:15]:
                doc.add_paragraph(f["text"], style="List Bullet")

        elif "entit" in sl:
            doc.add_heading("Extrahierte Entitaeten", level=1)
            if entities_data:
                for cat, items in entities_data.items():
                    if items:
                        doc.add_heading(cat.capitalize(), level=2)
                        for item in items[:10]:
                            val = item.get("value", item) if isinstance(item, dict) else item
                            doc.add_paragraph(str(val), style="List Bullet")

        elif "dokument" in sl or "uebersicht" in sl:
            doc.add_heading("Dokumentenuebersicht", level=1)
            # Tabelle
            table = doc.add_table(rows=1, cols=4)
            table.style = "Light Grid Accent 1"
            hdr = table.rows[0].cells
            hdr[0].text = "Titel"
            hdr[1].text = "Typ"
            hdr[2].text = "Seiten"
            hdr[3].text = "Status"
            for d in documents[:20]:
                row = table.add_row().cells
                row[0].text = d.get("title", "") or (d.title if hasattr(d, "title") else "")
                row[1].text = (d.get("file_type", "") or (d.file_type if hasattr(d, "file_type") else "")).upper()
                row[2].text = str(d.get("page_count", 0) or (d.page_count if hasattr(d, "page_count") else 0))
                row[3].text = d.get("status", "") or (d.status.value if hasattr(d, "status") and d.status else "")

        elif "quell" in sl or "source" in sl:
            doc.add_heading("Quellenverzeichnis", level=1)
            doc.add_paragraph(
                f"Alle Inhalte stammen ausschliesslich aus der {type_label} '{store_name}'. "
                f"Kein Weltwissen wurde verwendet."
            )
            for d in documents:
                title = d.get("title", "") or (d.title if hasattr(d, "title") else "")
                doc.add_paragraph(title, style="List Number")

    # Footer
    doc.add_paragraph("")
    p = doc.add_paragraph("Agentischer Document Store | Komm.ONE | On-Premise | DSGVO-konform")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in p.runs:
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(*CI_GRAY)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    logger.info(f"DOCX Export: {len(requested_sections)} Abschnitte fuer '{store_name}'")
    return buf.getvalue()


def export_pdf(
    store_name: str,
    store_type: str,
    summary: str,
    takeaways: list,
    entities_data: dict = None,
    facts: list = None,
    documents: list = None,
) -> bytes:
    """
    PDF-Export der Live-View (Zusammenfassung, Entitaeten, Fakten).
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()

    # Custom Styles
    styles.add(ParagraphStyle(name="CITitle", parent=styles["Title"], textColor=HexColor("#003A40"), fontSize=22, spaceAfter=12))
    styles.add(ParagraphStyle(name="CIHeading", parent=styles["Heading2"], textColor=HexColor("#00B2A9"), fontSize=14, spaceBefore=16, spaceAfter=6))
    styles.add(ParagraphStyle(name="CIBody", parent=styles["Normal"], fontSize=10, leading=14, spaceAfter=6))
    styles.add(ParagraphStyle(name="CISmall", parent=styles["Normal"], fontSize=8, textColor=HexColor("#6c757d")))

    story = []
    type_label = "Akte" if store_type == "akte" else "WissensDB"

    # Titel
    story.append(Paragraph(f"Live-View: {store_name}", styles["CITitle"]))
    story.append(Paragraph(f"{type_label} | Agentischer Document Store", styles["CISmall"]))
    story.append(Spacer(1, 16))

    # Zusammenfassung
    story.append(Paragraph("Zusammenfassung", styles["CIHeading"]))
    story.append(Paragraph(summary or "Keine Zusammenfassung verfuegbar.", styles["CIBody"]))
    story.append(Spacer(1, 8))

    # Kernfakten
    if takeaways:
        story.append(Paragraph("Kernfakten", styles["CIHeading"]))
        for tk in takeaways[:10]:
            val = tk.get("takeaway", tk.get("entity", ""))
            cnt = tk.get("count", 0)
            story.append(Paragraph(f"<bullet>&bull;</bullet> {val} ({cnt}x)", styles["CIBody"]))
        story.append(Spacer(1, 8))

    # Entitaeten
    if entities_data:
        story.append(Paragraph("Extrahierte Entitaeten", styles["CIHeading"]))
        for cat, items in entities_data.items():
            if items:
                story.append(Paragraph(f"<b>{cat.capitalize()}</b>", styles["CIBody"]))
                for item in items[:8]:
                    val = item.get("value", str(item)) if isinstance(item, dict) else str(item)
                    story.append(Paragraph(f"  - {val}", styles["CIBody"]))
        story.append(Spacer(1, 8))

    # Fakten
    if facts:
        story.append(Paragraph("Destillierte Fakten", styles["CIHeading"]))
        for f in facts[:10]:
            story.append(Paragraph(f"<bullet>&bull;</bullet> {f.get('text', '')[:150]}", styles["CIBody"]))
        story.append(Spacer(1, 8))

    # Dokumentenliste
    if documents:
        story.append(Paragraph("Dokumente", styles["CIHeading"]))
        table_data = [["Titel", "Typ", "Seiten"]]
        for d in documents[:20]:
            t = d.get("title", "") or (d.title if hasattr(d, "title") else "")
            ft = (d.get("file_type", "") or (d.file_type if hasattr(d, "file_type") else "")).upper()
            pg = str(d.get("page_count", 0) or (d.page_count if hasattr(d, "page_count") else 0))
            table_data.append([t[:50], ft, pg])

        t = Table(table_data, colWidths=[10*cm, 2*cm, 2*cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), HexColor("#003A40")),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#dee2e6")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [HexColor("#FFFFFF"), HexColor("#f2f5f5")]),
        ]))
        story.append(t)

    # Footer
    story.append(Spacer(1, 24))
    story.append(Paragraph(
        f"Alle Inhalte stammen ausschliesslich aus der {type_label} '{store_name}'. "
        "Agentischer Document Store | Komm.ONE | On-Premise | DSGVO-konform",
        styles["CISmall"]
    ))

    doc.build(story)
    buf.seek(0)
    logger.info(f"PDF Export: Live-View fuer '{store_name}'")
    return buf.getvalue()
