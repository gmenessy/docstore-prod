"""
Briefing-Export-Service — generiert PPTX/DOCX/PDF aus einem Decision-Briefing.

Nutzt die existierende Export-Infrastruktur (Komm.ONE-CI), aber spezifisch
zugeschnitten auf die 4-Fragen-Struktur des Briefings:
- Titel-Slide/Seite
- Sachstand
- Risiken (mit Schweregrad-Farbe)
- Naechste Schritte (als Tabelle)
- KI-Loesungsvorschlag

Fuer Gremiensitzungen, Protokoll-Vorlagen und Stakeholder-PDFs.
"""
import logging
import io
from datetime import datetime

logger = logging.getLogger(__name__)

# Komm.ONE CI
CI_MIDNIGHT = (0x00, 0x3A, 0x40)
CI_LAGOON = (0x00, 0xB2, 0xA9)
CI_AMARILLO = (0xF1, 0xC4, 0x00)
CI_BASIL = (0x00, 0x96, 0x5E)
CI_BURG = (0xE8, 0x96, 0x00)
CI_RED = (0xD9, 0x00, 0x00)
CI_WHITE = (0xFF, 0xFF, 0xFF)
CI_GRAY = (0x6C, 0x75, 0x7D)
CI_LIGHT = (0xF2, 0xF5, 0xF5)


def _severity_rgb(sev: str) -> tuple:
    return {
        "rot": CI_RED,
        "amber": CI_BURG,
        "gelb": CI_AMARILLO,
    }.get(sev, CI_GRAY)


def export_briefing_pptx(briefing: dict) -> bytes:
    """
    Gremien-taugliche PPTX: eine Folie pro Frage.
    Layout 16:9, Komm.ONE-CI, klare Typografie.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    store = briefing.get("store", {})
    sachstand = briefing.get("sachstand", {})
    risiken = briefing.get("risiken", {})
    next_steps = briefing.get("naechste_schritte", [])
    solution = briefing.get("loesungsvorschlag", {})

    blank_layout = prs.slide_layouts[6]

    def add_title_band(slide, title: str, subtitle: str = ""):
        """Komm.ONE-Titelband oben auf jeder Folie."""
        # Midnight-Band
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(*CI_MIDNIGHT)
        band.line.fill.background()

        # Lagoon-Akzent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), prs.slide_width, Inches(0.08)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*CI_LAGOON)
        accent.line.fill.background()

        # Titel-Text
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.7))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*CI_WHITE)
        p.font.name = "Arial"

        if subtitle:
            tbs = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12), Inches(0.35))
            tfs = tbs.text_frame
            tfs.margin_left = tfs.margin_right = tfs.margin_top = tfs.margin_bottom = 0
            ps = tfs.paragraphs[0]
            ps.text = subtitle
            ps.font.size = Pt(12)
            ps.font.color.rgb = RGBColor(180, 220, 218)
            ps.font.name = "Arial"

    def add_footer(slide, page_num: int, total: int):
        """Fusszeile mit Paginierung."""
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.3))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = f"Komm.ONE KI-Labor · Entscheider-Briefing · {page_num}/{total}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(*CI_GRAY)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.RIGHT

    # ── SLIDE 1: Titel ──
    slide = prs.slides.add_slide(blank_layout)
    # Vollflaechig Midnight
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*CI_MIDNIGHT)
    bg.line.fill.background()
    # Lagoon-Balken links
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3), Inches(0.4), Inches(1.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*CI_LAGOON)
    bar.line.fill.background()
    # Header-Label
    tb0 = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11), Inches(0.4))
    tf0 = tb0.text_frame; p0 = tf0.paragraphs[0]
    p0.text = "ENTSCHEIDER-BRIEFING"
    p0.font.size = Pt(12); p0.font.bold = True
    p0.font.color.rgb = RGBColor(*CI_LAGOON); p0.font.name = "Arial"
    # Titel
    tb1 = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.3))
    tf1 = tb1.text_frame; p1 = tf1.paragraphs[0]
    p1.text = store.get("name", "Sammlung")
    p1.font.size = Pt(36); p1.font.bold = True
    p1.font.color.rgb = RGBColor(*CI_WHITE); p1.font.name = "Arial"
    # Untertitel
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5))
    tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
    p2.text = f"{store.get('doc_count', 0)} Dokumente · {risiken.get('total', 0)} Risiken identifiziert · {len(next_steps)} Massnahmen"
    p2.font.size = Pt(14); p2.font.color.rgb = RGBColor(180, 220, 218); p2.font.name = "Arial"
    # Datum
    tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4))
    tf3 = tb3.text_frame; p3 = tf3.paragraphs[0]
    p3.text = f"Erstellt am {datetime.now().strftime('%d.%m.%Y')} · Komm.ONE KI-Labor"
    p3.font.size = Pt(10); p3.font.color.rgb = RGBColor(150, 170, 175); p3.font.name = "Arial"

    # ── SLIDE 2: Sachstand ──
    slide = prs.slides.add_slide(blank_layout)
    add_title_band(slide, "1  ·  Sachstand", f"Synthese aus {sachstand.get('sources', 0)} Dokumenten")
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = sachstand.get("text", "–")
    p.font.size = Pt(18); p.font.color.rgb = RGBColor(*CI_MIDNIGHT); p.font.name = "Arial"
    # Meta
    tbm = slide.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12), Inches(0.4))
    tfm = tbm.text_frame; pm = tfm.paragraphs[0]
    model = sachstand.get("model", "")
    conf = int((sachstand.get("confidence") or 0) * 100)
    pm.text = f"Confidence {conf}% · {'KI-generiert' if model == 'llm' else 'Extraktiv'}"
    pm.font.size = Pt(10); pm.font.color.rgb = RGBColor(*CI_GRAY); pm.font.name = "Consolas"
    add_footer(slide, 2, 5)

    # ── SLIDE 3: Risiken ──
    slide = prs.slides.add_slide(blank_layout)
    add_title_band(slide, "2  ·  Risiken", f"{risiken.get('total', 0)} identifiziert")
    risks_list = risiken.get("risks", [])[:6]
    y_offset = 1.5
    for r in risks_list:
        # Farb-Bar links
        sev_rgb = _severity_rgb(r.get("severity", "gelb"))
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y_offset), Inches(0.08), Inches(0.7)
        )
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*sev_rgb)
        bar.line.fill.background()
        # Titel
        tb = slide.shapes.add_textbox(Inches(0.85), Inches(y_offset - 0.02), Inches(11.5), Inches(0.4))
        tf = tb.text_frame; p = tf.paragraphs[0]
        p.text = r.get("title", "")
        p.font.size = Pt(13); p.font.bold = True
        p.font.color.rgb = RGBColor(*CI_MIDNIGHT); p.font.name = "Arial"
        # Beschreibung
        tbd = slide.shapes.add_textbox(Inches(0.85), Inches(y_offset + 0.3), Inches(11.5), Inches(0.4))
        tfd = tbd.text_frame; pd = tfd.paragraphs[0]
        desc = (r.get("description") or "")[:160].strip()
        pd.text = desc + (" …" if r.get("description", "") and len(r["description"]) > 160 else "")
        pd.font.size = Pt(10); pd.font.color.rgb = RGBColor(*CI_GRAY); pd.font.name = "Arial"
        y_offset += 0.82
    add_footer(slide, 3, 5)

    # ── SLIDE 4: Naechste Schritte ──
    slide = prs.slides.add_slide(blank_layout)
    add_title_band(slide, "3  ·  Naechste Schritte", f"{len(next_steps)} Massnahmen priorisiert")
    if next_steps:
        # Tabelle
        rows = min(len(next_steps), 6) + 1
        cols = 4
        left = Inches(0.6); top = Inches(1.6)
        width = Inches(12.1); height = Inches(0.4 * rows + 0.2)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        tbl = table_shape.table
        # Kopfzeile
        headers = ["#", "Massnahme", "Zustaendig", "Frist"]
        col_widths = [Inches(0.5), Inches(7.3), Inches(2.1), Inches(2.2)]
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(*CI_MIDNIGHT)
            cell.text = h
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11); para.font.bold = True
                para.font.color.rgb = RGBColor(*CI_WHITE); para.font.name = "Arial"
        # Datenzeilen
        for idx, step in enumerate(next_steps[:6], start=1):
            cells = [
                str(idx),
                step.get("title", ""),
                step.get("assignee", "–") or "–",
                step.get("due_date", "–") or "–",
            ]
            for ci, text in enumerate(cells):
                cell = tbl.cell(idx, ci)
                cell.text = str(text)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(10)
                    para.font.color.rgb = RGBColor(*CI_MIDNIGHT); para.font.name = "Arial"
    else:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(12), Inches(1))
        tf = tb.text_frame; p = tf.paragraphs[0]
        p.text = "Keine offenen Massnahmen."
        p.font.size = Pt(14); p.font.color.rgb = RGBColor(*CI_GRAY); p.font.name = "Arial"
    add_footer(slide, 4, 5)

    # ── SLIDE 5: KI-Loesungsvorschlag ──
    slide = prs.slides.add_slide(blank_layout)
    add_title_band(slide, "4  ·  KI-Loesungsvorschlag", "Synthese aus Sachstand, Risiken und Massnahmen")
    # Lagoon-akzentuierte Karte
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.5)
    )
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(224, 244, 243)
    card.line.fill.background()
    # Lagoon-Balken
    lb = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.6), Inches(0.08), Inches(4.5)
    )
    lb.fill.solid(); lb.fill.fore_color.rgb = RGBColor(*CI_LAGOON); lb.line.fill.background()
    # Label
    tbl_ = slide.shapes.add_textbox(Inches(0.9), Inches(1.75), Inches(11), Inches(0.4))
    tf_ = tbl_.text_frame; pl = tf_.paragraphs[0]
    pl.text = "EMPFEHLUNG"
    pl.font.size = Pt(10); pl.font.bold = True
    pl.font.color.rgb = RGBColor(0, 120, 115); pl.font.name = "Arial"
    # Text
    tb2_ = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(3.5))
    tf2_ = tb2_.text_frame; tf2_.word_wrap = True
    p2_ = tf2_.paragraphs[0]
    p2_.text = solution.get("text", "–")
    p2_.font.size = Pt(16); p2_.font.color.rgb = RGBColor(*CI_MIDNIGHT); p2_.font.name = "Arial"
    # Meta
    model_s = solution.get("model", "–")
    sources = solution.get("sources", 0)
    conf_s = int((solution.get("confidence") or 0) * 100)
    tbm_ = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(12), Inches(0.3))
    pm_ = tbm_.text_frame.paragraphs[0]
    pm_.text = f"Generiert mit {model_s} · {sources} Quellen · Confidence {conf_s}%"
    pm_.font.size = Pt(9); pm_.font.color.rgb = RGBColor(*CI_GRAY); pm_.font.name = "Consolas"
    add_footer(slide, 5, 5)

    # In Bytes schreiben
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def export_briefing_docx(briefing: dict) -> bytes:
    """
    Protokoll-Vorlage als Word-Dokument.
    Strukturiert wie ein Tagesordnungspunkt.
    """
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_ALIGN_VERTICAL
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    store = briefing.get("store", {})
    sachstand = briefing.get("sachstand", {})
    risiken = briefing.get("risiken", {})
    next_steps = briefing.get("naechste_schritte", [])
    solution = briefing.get("loesungsvorschlag", {})

    # Standardschrift
    styles = doc.styles
    styles["Normal"].font.name = "Arial"
    styles["Normal"].font.size = Pt(11)

    # Kopfzeile
    section = doc.sections[0]
    section.top_margin = Cm(2); section.bottom_margin = Cm(2)
    section.left_margin = Cm(2); section.right_margin = Cm(2)

    # Titel-Block
    title_p = doc.add_paragraph()
    title_run = title_p.add_run("ENTSCHEIDER-BRIEFING")
    title_run.bold = True; title_run.font.size = Pt(10)
    title_run.font.color.rgb = RGBColor(0x00, 0xB2, 0xA9)

    h1 = doc.add_paragraph()
    h1_run = h1.add_run(store.get("name", "Sammlung"))
    h1_run.bold = True; h1_run.font.size = Pt(18)
    h1_run.font.color.rgb = RGBColor(0x00, 0x3A, 0x40)

    meta_p = doc.add_paragraph()
    meta_run = meta_p.add_run(
        f"Erstellt am {datetime.now().strftime('%d.%m.%Y')} · "
        f"{store.get('doc_count', 0)} Dokumente · "
        f"{risiken.get('total', 0)} Risiken · "
        f"{len(next_steps)} Massnahmen"
    )
    meta_run.font.size = Pt(10); meta_run.font.color.rgb = RGBColor(0x6C, 0x75, 0x7D)

    doc.add_paragraph()

    # ── 1. Sachstand ──
    h = doc.add_paragraph()
    r = h.add_run("1.  Sachstand")
    r.bold = True; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x00, 0x3A, 0x40)

    p = doc.add_paragraph(sachstand.get("text", "–"))
    for run in p.runs:
        run.font.size = Pt(11)

    m = doc.add_paragraph()
    mr = m.add_run(
        f"  Synthese aus {sachstand.get('sources', 0)} Dokumenten · "
        f"Confidence {int((sachstand.get('confidence') or 0) * 100)}% · "
        f"{'KI-generiert' if sachstand.get('model') == 'llm' else 'Extraktiv'}"
    )
    mr.italic = True; mr.font.size = Pt(9); mr.font.color.rgb = RGBColor(0x6C, 0x75, 0x7D)

    doc.add_paragraph()

    # ── 2. Risiken ──
    h = doc.add_paragraph()
    r = h.add_run("2.  Risiken")
    r.bold = True; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x00, 0x3A, 0x40)

    risks_list = risiken.get("risks", [])
    if not risks_list:
        doc.add_paragraph("Keine Risiken identifiziert.")
    else:
        for risk in risks_list[:8]:
            sev = risk.get("severity", "gelb")
            sev_label = {"rot": "AKUT", "amber": "BEOBACHTEN", "gelb": "HINWEIS"}.get(sev, sev.upper())
            sev_color = {"rot": RGBColor(0xD9, 0x00, 0x00), "amber": RGBColor(0xE8, 0x96, 0x00), "gelb": RGBColor(0xF1, 0xC4, 0x00)}.get(sev, RGBColor(0x6C, 0x75, 0x7D))
            rp = doc.add_paragraph()
            s = rp.add_run(f"[{sev_label}]  ")
            s.bold = True; s.font.size = Pt(9); s.font.color.rgb = sev_color
            t = rp.add_run(risk.get("title", ""))
            t.bold = True; t.font.size = Pt(11)
            if risk.get("description"):
                dp = doc.add_paragraph()
                dr = dp.add_run("  " + risk["description"][:300])
                dr.font.size = Pt(10); dr.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    doc.add_paragraph()

    # ── 3. Naechste Schritte ──
    h = doc.add_paragraph()
    r = h.add_run("3.  Naechste Schritte")
    r.bold = True; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x00, 0x3A, 0x40)

    if not next_steps:
        doc.add_paragraph("Keine offenen Massnahmen.")
    else:
        # Tabelle
        tbl = doc.add_table(rows=1 + min(len(next_steps), 10), cols=5)
        tbl.style = "Light Grid Accent 1"
        hdr = tbl.rows[0].cells
        headers = ["#", "Massnahme", "Prio", "Zustaendig", "Frist"]
        for i, h_ in enumerate(headers):
            cell_p = hdr[i].paragraphs[0]
            run = cell_p.add_run(h_)
            run.bold = True; run.font.size = Pt(10)
        for idx, step in enumerate(next_steps[:10], start=1):
            row = tbl.rows[idx].cells
            row[0].text = str(idx)
            row[1].text = step.get("title", "")
            row[2].text = step.get("priority", "").capitalize()
            row[3].text = step.get("assignee", "–") or "–"
            row[4].text = step.get("due_date", "–") or "–"
            for c in row:
                for par in c.paragraphs:
                    for run in par.runs:
                        run.font.size = Pt(10)

    doc.add_paragraph()

    # ── 4. KI-Loesungsvorschlag ──
    h = doc.add_paragraph()
    r = h.add_run("4.  KI-Loesungsvorschlag")
    r.bold = True; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x00, 0x3A, 0x40)

    # Empfehlungs-Karte
    ep = doc.add_paragraph()
    eh = ep.add_run("EMPFEHLUNG")
    eh.bold = True; eh.font.size = Pt(9); eh.font.color.rgb = RGBColor(0x00, 0x78, 0x73)

    body_p = doc.add_paragraph(solution.get("text", "–"))
    body_p.paragraph_format.left_indent = Cm(0.3)
    for run in body_p.runs:
        run.font.size = Pt(11)

    meta2 = doc.add_paragraph()
    mr2 = meta2.add_run(
        f"  Generiert mit {solution.get('model', '–')} · "
        f"{solution.get('sources', 0)} Quellen · "
        f"Confidence {int((solution.get('confidence') or 0) * 100)}%"
    )
    mr2.italic = True; mr2.font.size = Pt(9); mr2.font.color.rgb = RGBColor(0x6C, 0x75, 0x7D)

    doc.add_paragraph()
    footer = doc.add_paragraph()
    fr = footer.add_run("Komm.ONE KI-Labor · Entscheider-Briefing · Automatisch generiert")
    fr.italic = True; fr.font.size = Pt(9); fr.font.color.rgb = RGBColor(0x6C, 0x75, 0x7D)
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def export_briefing_pdf(briefing: dict) -> bytes:
    """
    1-Seiten PDF-Briefing fuer Mail-Versand / Schnell-Ausdruck.
    Kompakter als die Word-Variante.
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, KeepTogether,
    )
    from reportlab.lib.enums import TA_LEFT

    store = briefing.get("store", {})
    sachstand = briefing.get("sachstand", {})
    risiken = briefing.get("risiken", {})
    next_steps = briefing.get("naechste_schritte", [])
    solution = briefing.get("loesungsvorschlag", {})

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=1.5 * cm, rightMargin=1.5 * cm,
        topMargin=1.5 * cm, bottomMargin=1.5 * cm,
    )

    styles = getSampleStyleSheet()
    midnight_hex = colors.HexColor("#003A40")
    lagoon_hex = colors.HexColor("#00B2A9")
    gray_hex = colors.HexColor("#6C757D")

    title_style = ParagraphStyle(
        "Title", parent=styles["Heading1"], fontSize=18, fontName="Helvetica-Bold",
        textColor=midnight_hex, spaceAfter=4, spaceBefore=0,
    )
    label_style = ParagraphStyle(
        "Label", parent=styles["Normal"], fontSize=8, fontName="Helvetica-Bold",
        textColor=lagoon_hex, spaceAfter=2,
    )
    section_style = ParagraphStyle(
        "Section", parent=styles["Heading2"], fontSize=12, fontName="Helvetica-Bold",
        textColor=midnight_hex, spaceAfter=3, spaceBefore=10,
    )
    body_style = ParagraphStyle(
        "Body", parent=styles["Normal"], fontSize=10, fontName="Helvetica",
        textColor=midnight_hex, spaceAfter=4, leading=13,
    )
    meta_style = ParagraphStyle(
        "Meta", parent=styles["Normal"], fontSize=8, fontName="Helvetica-Oblique",
        textColor=gray_hex, spaceAfter=8,
    )

    story = []

    # Header
    story.append(Paragraph("ENTSCHEIDER-BRIEFING", label_style))
    story.append(Paragraph(store.get("name", "Sammlung"), title_style))
    meta_text = (
        f"Erstellt {datetime.now().strftime('%d.%m.%Y')} &middot; "
        f"{store.get('doc_count', 0)} Dokumente &middot; "
        f"{risiken.get('total', 0)} Risiken &middot; "
        f"{len(next_steps)} Massnahmen"
    )
    story.append(Paragraph(meta_text, meta_style))

    # Trennlinie
    hr = Table([[""]], colWidths=[18 * cm], rowHeights=[1])
    hr.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), lagoon_hex)]))
    story.append(hr)
    story.append(Spacer(1, 8))

    # 1. Sachstand
    story.append(Paragraph("1. Sachstand", section_style))
    story.append(Paragraph(sachstand.get("text", "–"), body_style))
    story.append(Paragraph(
        f"Synthese aus {sachstand.get('sources', 0)} Dokumenten &middot; "
        f"Confidence {int((sachstand.get('confidence') or 0) * 100)}%",
        meta_style
    ))

    # 2. Risiken
    story.append(Paragraph("2. Risiken", section_style))
    risks_list = risiken.get("risks", [])[:5]
    if not risks_list:
        story.append(Paragraph("Keine Risiken identifiziert.", body_style))
    else:
        rows = []
        for r in risks_list:
            sev = r.get("severity", "gelb")
            sev_label = {"rot": "AKUT", "amber": "BEOB.", "gelb": "HINWEIS"}.get(sev, "–")
            rows.append([sev_label, r.get("title", "")])
        risk_tbl = Table(rows, colWidths=[2.5 * cm, 15.5 * cm])
        style_list = [
            ("FONT", (0, 0), (-1, -1), "Helvetica", 9),
            ("TEXTCOLOR", (1, 0), (1, -1), midnight_hex),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("FONT", (0, 0), (0, -1), "Helvetica-Bold", 8),
        ]
        for i, r in enumerate(risks_list):
            sev = r.get("severity", "gelb")
            c = {"rot": colors.HexColor("#D90000"), "amber": colors.HexColor("#E89600"), "gelb": colors.HexColor("#F1C400")}.get(sev, gray_hex)
            style_list.append(("TEXTCOLOR", (0, i), (0, i), c))
        risk_tbl.setStyle(TableStyle(style_list))
        story.append(risk_tbl)

    # 3. Naechste Schritte
    story.append(Paragraph("3. Naechste Schritte", section_style))
    if not next_steps:
        story.append(Paragraph("Keine offenen Massnahmen.", body_style))
    else:
        header = [["#", "Massnahme", "Prio", "Zustaendig", "Frist"]]
        data_rows = []
        for idx, step in enumerate(next_steps[:6], start=1):
            data_rows.append([
                str(idx),
                step.get("title", "")[:60],
                (step.get("priority", "") or "").capitalize(),
                (step.get("assignee", "–") or "–")[:20],
                step.get("due_date", "–") or "–",
            ])
        step_tbl = Table(header + data_rows, colWidths=[0.7 * cm, 9 * cm, 2 * cm, 3.3 * cm, 3 * cm])
        step_tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), midnight_hex),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONT", (0, 0), (-1, 0), "Helvetica-Bold", 9),
            ("FONT", (0, 1), (-1, -1), "Helvetica", 9),
            ("TEXTCOLOR", (0, 1), (-1, -1), midnight_hex),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F5F7F7")]),
        ]))
        story.append(step_tbl)

    # 4. KI-Loesungsvorschlag
    story.append(Paragraph("4. KI-Loesungsvorschlag", section_style))
    sol_inner = Paragraph(solution.get("text", "–"), body_style)
    sol_tbl = Table([[sol_inner]], colWidths=[17 * cm])
    sol_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#E1F5EE")),
        ("LINEBEFORE", (0, 0), (0, -1), 3, lagoon_hex),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("LEFTPADDING", (0, 0), (-1, -1), 12),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
    ]))
    story.append(sol_tbl)
    story.append(Paragraph(
        f"Generiert mit {solution.get('model', '–')} &middot; "
        f"{solution.get('sources', 0)} Quellen &middot; "
        f"Confidence {int((solution.get('confidence') or 0) * 100)}%",
        meta_style
    ))

    # Fusszeile
    story.append(Spacer(1, 12))
    story.append(Paragraph("Komm.ONE KI-Labor · Automatisch generiert", meta_style))

    doc.build(story)
    return buf.getvalue()
